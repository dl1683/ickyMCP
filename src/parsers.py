"""Document parsers for extracting text from various file types."""

from pathlib import Path
from typing import Optional
from dataclasses import dataclass

from .config import SUPPORTED_EXTENSIONS


@dataclass
class ParsedDocument:
    """Result of parsing a document."""
    text: str
    page_count: int
    file_type: str
    metadata: dict


def parse_text_file(path: Path) -> ParsedDocument:
    """Parse plain text or markdown file."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    return ParsedDocument(
        text=text,
        page_count=1,
        file_type="text" if path.suffix == ".txt" else "markdown",
        metadata={}
    )


def parse_pdf(path: Path) -> ParsedDocument:
    """Parse PDF file using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []

    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            # Add page marker for tracking
            pages.append(f"[PAGE {page_num}]\n{text}")

    # Extract metadata safely
    pdf_meta = {}
    if reader.metadata:
        try:
            for key in ['/Title', '/Author', '/Subject', '/Creator']:
                if key in reader.metadata:
                    pdf_meta[key.strip('/')] = str(reader.metadata[key])
        except Exception:
            pass

    return ParsedDocument(
        text="\n\n".join(pages),
        page_count=len(reader.pages),
        file_type="pdf",
        metadata={"pdf_metadata": pdf_meta}
    )


def parse_word(path: Path) -> ParsedDocument:
    """Parse Word document (.docx) using python-docx.

    Adds estimated page markers based on character count (~3000 chars per page).
    """
    from docx import Document

    doc = Document(str(path))
    paragraphs = []

    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    # Join paragraphs
    full_text = "\n\n".join(paragraphs)

    # Estimate page count and add page markers (~3000 chars per page)
    CHARS_PER_PAGE = 3000
    total_chars = len(full_text)
    estimated_pages = max(1, (total_chars + CHARS_PER_PAGE - 1) // CHARS_PER_PAGE)

    # Add page markers at appropriate positions
    if total_chars > CHARS_PER_PAGE:
        marked_text = []
        current_page = 1
        current_pos = 0

        marked_text.append(f"[PAGE {current_page}]")

        for para in paragraphs:
            marked_text.append(para)
            current_pos += len(para) + 2  # +2 for \n\n

            # Check if we've crossed a page boundary
            expected_page = (current_pos // CHARS_PER_PAGE) + 1
            if expected_page > current_page and expected_page <= estimated_pages:
                current_page = expected_page
                marked_text.append(f"\n[PAGE {current_page}]")

        full_text = "\n\n".join(marked_text)
    else:
        # Small doc - just add single page marker
        full_text = f"[PAGE 1]\n{full_text}"

    return ParsedDocument(
        text=full_text,
        page_count=estimated_pages,
        file_type="word",
        metadata={"estimated_pages": True}
    )


def parse_powerpoint(path: Path) -> ParsedDocument:
    """Parse PowerPoint presentation (.pptx) using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_text.append(f"[SLIDE {slide_num}]")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

            # Handle tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_text.append(row_text)

        if len(slide_text) > 1:  # More than just the slide marker
            slides.append("\n".join(slide_text))

    return ParsedDocument(
        text="\n\n".join(slides),
        page_count=len(prs.slides),
        file_type="powerpoint",
        metadata={}
    )


def parse_excel(path: Path) -> ParsedDocument:
    """Parse Excel spreadsheet (.xlsx) using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = [f"[SHEET: {sheet_name}]"]

        for row in sheet.iter_rows():
            row_values = []
            for cell in row:
                if cell.value is not None:
                    row_values.append(str(cell.value))

            if row_values:
                sheet_text.append(" | ".join(row_values))

        if len(sheet_text) > 1:  # More than just the sheet marker
            sheets.append("\n".join(sheet_text))

    return ParsedDocument(
        text="\n\n".join(sheets),
        page_count=len(wb.sheetnames),
        file_type="excel",
        metadata={"sheet_names": wb.sheetnames}
    )


# Parser registry
PARSERS = {
    ".txt": parse_text_file,
    ".md": parse_text_file,
    ".pdf": parse_pdf,
    ".docx": parse_word,
    ".pptx": parse_powerpoint,
    ".xlsx": parse_excel,
}


def parse_document(path: Path) -> Optional[ParsedDocument]:
    """
    Parse a document based on its file extension.

    Returns None if the file type is not supported.
    """
    suffix = path.suffix.lower()

    if suffix not in PARSERS:
        return None

    parser = PARSERS[suffix]
    return parser(path)


def is_supported(path: Path) -> bool:
    """Check if a file type is supported."""
    return path.suffix.lower() in SUPPORTED_EXTENSIONS
